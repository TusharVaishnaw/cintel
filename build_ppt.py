#!/usr/bin/env python3
"""
build_ppt.py <slide_data.json> <out.pptx>

Python port of build_ppt.js — same data-driven approach: every add_* function
skips itself or its rows when the underlying field is missing. Never prints a
placeholder for null data. Drop-in replacement for the Node/pptxgenjs stage;
takes the same slide_data.json produced by extract_slide_data.py.

Usage:
    python build_ppt.py slide_data.json out.pptx

Install:
    pip install python-pptx --break-system-packages
"""
import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn

FONT = "Calibri"
NAVY = RGBColor(0x1F, 0x2D, 0x50)
NAVY_DK = RGBColor(0x13, 0x19, 0x2E)
ACCENT = RGBColor(0x2E, 0x6F, 0x9E)
SLATE = RGBColor(0x5B, 0x64, 0x72)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOOD = RGBColor(0x3C, 0x8A, 0x5B)
WARN = RGBColor(0xC6, 0x86, 0x2B)
BAD = RGBColor(0xB5, 0x47, 0x3A)
LINE_GREY = RGBColor(0xE3, 0xE7, 0xEE)

PW, PH = 13.33, 7.5  # inches, matches pptxgenjs LAYOUT_WIDE


def IN(v):
    return Inches(v)


# ---------------------------------------------------------------------------
# low-level helpers
# ---------------------------------------------------------------------------
def new_slide(prs, bg=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    if bg is not None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg
    return slide


def add_textbox(slide, x, y, w, h, text, *, size=11, bold=False, color=NAVY_DK,
                 align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT,
                 hyperlink=None, wrap=True, shrink=False):
    """Single-run textbox. shrink=True enables autofit-shrink-to-fit (best
    equivalent of pptxgenjs's fit:'shrink')."""
    box = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = valign
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if shrink:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if hyperlink:
        run.hyperlink.address = hyperlink
    return box


def add_bullets(slide, x, y, w, h, items, *, size=10, color=NAVY_DK, font=FONT,
                 valign=MSO_ANCHOR.TOP):
    """items: list of plain strings, each its own bulleted paragraph."""
    box = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_bullet(p)
        run = p.add_run()
        run.text = str(item)
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def _set_bullet(paragraph):
    """Turn on a real bullet char via raw XML (python-pptx has no bullet API)."""
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '•'})
    pPr.append(buChar)
    marL = pPr.makeelement(qn('a:marL'), {})
    pPr.set('marL', '182880')  # ~0.2in indent so bullet text doesn't hug the glyph
    pPr.set('indent', '-182880')


def add_card(slide, x, y, w, h, *, fill=CARD_BG, line=LINE_GREY):
    """Rounded rect card. python-pptx shadow control is limited to on/off —
    inherits the shape's default (subtle) shadow; explicit soft-shadow tuning
    like pptxgenjs's blur/offset isn't exposed without raw XML, so a light
    default shadow is left on for the same "lifted card" visual."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(x), IN(y), IN(w), IN(h))
    shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_rect(slide, x, y, w, h, fill, line_none=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(x), IN(y), IN(w), IN(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_none:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_ellipse(slide, x, y, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, IN(x), IN(y), IN(w), IN(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_circle_image(slide, path, x, y, w, h):
    """Crop an image to a circle — python-pptx exposes the oval-crop geometry
    via the picture shape's auto_shape_type equivalent by setting the
    preset geometry directly on the picture's XML."""
    pic = slide.shapes.add_picture(path, IN(x), IN(y), IN(w), IN(h))
    pic.line.fill.background()
    sp = pic._element
    spPr = sp.find(qn('p:spPr'))
    geom = spPr.makeelement(qn('a:prstGeom'), {'prst': 'ellipse'})
    av_lst = geom.makeelement(qn('a:avLst'), {})
    geom.append(av_lst)
    # remove any existing geometry node before inserting the ellipse crop
    for tag in ('a:prstGeom', 'a:custGeom'):
        existing = spPr.find(qn(tag))
        if existing is not None:
            spPr.remove(existing)
    spPr.append(geom)
    return pic


def footer(slide, company_name, label=None):
    add_textbox(slide, 0.4, PH - 0.35, 8, 0.3, label or company_name, size=9, color=SLATE)
    add_textbox(slide, PW - 3.4, PH - 0.35, 3.0, 0.3, "Company Intelligence",
                size=9, color=SLATE, align=PP_ALIGN.RIGHT)


def slide_title(slide, text):
    add_textbox(slide, 0.5, 0.35, PW - 1.0, 0.65, text, size=26, bold=True, color=NAVY)


def money_or_none(v):
    return v if v not in (None, "", [], {}) else None


# ---------------------------------------------------------------------------
# Slide builders — one function per slide type, mirroring build_ppt.js 1:1
# ---------------------------------------------------------------------------
def add_title_slide(prs, d):
    sl = new_slide(prs, NAVY_DK)
    add_textbox(sl, 0.8, 2.7, PW - 1.6, 1.3, d["company_name"], size=44, bold=True, color=WHITE)
    add_textbox(sl, 0.8, 3.9, PW - 1.6, 0.6, "Company Intelligence", size=20, color=RGBColor(0xAF, 0xC6, 0xE3))
    if d.get("ticker"):
        add_textbox(sl, 0.8, 4.45, 4, 0.4, d["ticker"], size=13, color=RGBColor(0x7E, 0x96, 0xBF))
    add_textbox(sl, 0.8, PH - 0.7, 6, 0.4, "Information Resource Centre", size=11, color=RGBColor(0x7E, 0x96, 0xBF))


def add_separator(prs, title):
    sl = new_slide(prs, NAVY)
    add_textbox(sl, 0.8, PH / 2 - 0.5, PW - 1.6, 1.0, title, size=32, bold=True, color=WHITE)


def add_overview(prs, d):
    sl = new_slide(prs)
    slide_title(sl, "Company Overview")
    ov = d["overview"]
    stats = [("Founded", ov.get("founded")), ("Headquarters", ov.get("headquarters")),
             ("Employees", ov.get("employees")), ("Revenue", ov.get("revenue")),
             ("Market Cap", ov.get("market_cap"))]
    stats = [(l, v) for l, v in stats if v]

    n = max(len(stats), 1)
    col_w = (PW - 1.0 - 0.3 * (len(stats) - 1)) / n if stats else 0
    for i, (label, val) in enumerate(stats):
        x = 0.5 + i * (col_w + 0.3)
        add_card(sl, x, 1.25, col_w, 1.1)
        add_textbox(sl, x + 0.12, 1.35, col_w - 0.24, 0.3, label.upper(), size=9, bold=True, color=SLATE)
        add_textbox(sl, x + 0.12, 1.62, col_w - 0.24, 0.65, str(val), size=12, bold=True, color=NAVY, shrink=True)

    y = 2.65
    lists = [("Key Competitors", ov.get("key_competitors") or []),
             ("Key Acquisitions", ov.get("key_acquisitions") or [])]
    lists = [(l, arr) for l, arr in lists if arr]
    two_col = [(l, v) for l, v in [("Brands", ov.get("brands")), ("Services", ov.get("services"))] if v]
    two_col_h = 1.55 if lists else (PH - y - 0.55)
    half_w = (PW - 1.0 - 0.3) / 2
    for i, (label, val) in enumerate(two_col):
        x = 0.5 + i * (half_w + 0.3)
        add_card(sl, x, y, half_w, two_col_h)
        add_textbox(sl, x + 0.15, y + 0.1, half_w - 0.3, 0.3, label, size=12, bold=True, color=ACCENT)
        add_textbox(sl, x + 0.15, y + 0.42, half_w - 0.3, two_col_h - 0.55, val, size=10, color=NAVY_DK, shrink=True)
    y += two_col_h + 0.25

    if lists:
        w2 = (PW - 1.0 - 0.3) / len(lists)
        for i, (label, arr) in enumerate(lists):
            x = 0.5 + i * (w2 + 0.3)
            h = PH - y - 0.55
            add_card(sl, x, y, w2, h)
            add_textbox(sl, x + 0.15, y + 0.1, w2 - 0.3, 0.3, label, size=12, bold=True, color=ACCENT)
            add_bullets(sl, x + 0.15, y + 0.42, w2 - 0.3, h - 0.55, arr, size=10, color=NAVY_DK)
    footer(sl, d["company_name"])


def add_mission_vision(prs, d):
    mv = d["mission_vision"]
    items = [(l, v) for l, v in [("Mission", mv.get("mission")), ("Vision", mv.get("vision")),
                                   ("Values", mv.get("values"))] if v]
    if not items:
        return
    sl = new_slide(prs)
    slide_title(sl, "Mission / Vision")
    h = (PH - 1.4 - 0.3 * (len(items) - 1)) / len(items)
    for i, (label, val) in enumerate(items):
        y = 1.25 + i * (h + 0.3)
        add_card(sl, 0.5, y, PW - 1.0, h)
        add_textbox(sl, 0.75, y + 0.12, 2.2, h - 0.24, label, size=15, bold=True, color=ACCENT, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(sl, 3.0, y + 0.12, PW - 3.5, h - 0.24, val, size=11.5, color=NAVY_DK, valign=MSO_ANCHOR.MIDDLE, shrink=True)
    footer(sl, d["company_name"])


def add_geo(prs, d):
    g = d["geo"]
    rows = [(l, v) for l, v in [("Countries", g.get("countries")), ("Regions", g.get("regions")),
                                  ("Offices / Facilities", g.get("offices")),
                                  ("Delivery Centers", g.get("delivery_centers")),
                                  ("Revenue by Geography", g.get("geographic_revenue"))] if v]
    if not rows:
        return
    sl = new_slide(prs)
    slide_title(sl, "Geographic Presence")
    y = 1.3
    for label, val in rows:
        val_str = str(val)
        h = 0.55 + (len(val_str) // 90) * 0.22
        add_card(sl, 0.5, y, PW - 1.0, h)
        add_textbox(sl, 0.7, y + 0.08, 2.6, h - 0.16, label, size=12, bold=True, color=ACCENT)
        add_textbox(sl, 3.3, y + 0.08, PW - 4.0, h - 0.16, val_str, size=10.5, color=NAVY_DK)
        y += h + 0.18
    footer(sl, d["company_name"])


def add_segments(prs, d):
    segs = d["segments"]
    if not segs:
        return
    sl = new_slide(prs)
    slide_title(sl, "Business Segments")
    n = len(segs)
    col_w = (PW - 1.0 - 0.25 * (n - 1)) / n
    for i, seg in enumerate(segs):
        x = 0.5 + i * (col_w + 0.25)
        h = 3.0
        add_card(sl, x, 1.25, col_w, h)
        add_textbox(sl, x + 0.12, 1.35, col_w - 0.24, 0.45, seg.get("name") or "Segment", size=12.5, bold=True, color=NAVY, shrink=True)
        iy = 1.85
        for lbl, val in [("Revenue", seg.get("revenue")), ("Growth", seg.get("growth")), ("Op. Margin", seg.get("op_margin"))]:
            if not val:
                continue
            add_textbox(sl, x + 0.12, iy, col_w - 0.24, 0.28, f"{lbl}: ", size=9.5, bold=True, color=SLATE)
            add_textbox(sl, x + 0.12, iy + 0.18, col_w - 0.24, 0.28, val, size=11.5, bold=True, color=NAVY_DK)
            iy += 0.5
        if seg.get("description"):
            add_textbox(sl, x + 0.12, iy + 0.05, col_w - 0.24, h - (iy + 0.05 - 1.25) - 0.1,
                        seg["description"], size=8.5, color=SLATE, shrink=True)

    with_pct = [s for s in segs if s.get("revenue_pct") is not None]
    if with_pct:
        add_textbox(sl, 0.5, 4.5, 4, 0.35, "Revenue Mix", size=13, bold=True, color=NAVY)
        chart_data = CategoryChartData()
        chart_data.categories = [s.get("name") or "Segment" for s in with_pct]
        chart_data.add_series("Revenue %", [s["revenue_pct"] for s in with_pct])
        gframe = sl.shapes.add_chart(XL_CHART_TYPE.PIE, IN(0.5), IN(4.85), IN(6.0), IN(2.35), chart_data)
        chart = gframe.chart
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.number_format = '0"%"'
        dl.number_format_is_linked = False
        dl.font.color.rgb = WHITE
        palette = [ACCENT, RGBColor(0x6F, 0xA8, 0xC9), NAVY, RGBColor(0x9D, 0xB4, 0xC9), SLATE]
        for i, point in enumerate(plot.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = palette[i % len(palette)]
    footer(sl, d["company_name"])


def add_sustainability_strategy(prs, d):
    s = d["sustainability"]
    if not s.get("strategy") and not s.get("goals"):
        return
    sl = new_slide(prs)
    slide_title(sl, "Sustainability Strategy and Goals")
    items = [(l, v) for l, v in [("Strategy", s.get("strategy")), ("Goals", s.get("goals"))] if v]
    h = (PH - 1.4 - 0.3 * (len(items) - 1)) / max(len(items), 1)
    for i, (label, val) in enumerate(items):
        y = 1.25 + i * (h + 0.3)
        add_card(sl, 0.5, y, PW - 1.0, h)
        add_textbox(sl, 0.75, y + 0.15, 2.5, 0.4, label, size=14, bold=True, color=ACCENT)
        add_textbox(sl, 0.75, y + 0.55, PW - 1.5, h - 0.7, val, size=10.5, color=NAVY_DK, shrink=True)
    footer(sl, d["company_name"])


def add_sustainability_initiatives(prs, d):
    s = d["sustainability"]
    if not s.get("key_initiatives"):
        return
    sl = new_slide(prs)
    slide_title(sl, "Sustainability – Key Initiatives")
    add_card(sl, 0.5, 1.25, PW - 1.0, PH - 1.85)
    add_textbox(sl, 0.8, 1.4, PW - 1.6, PH - 2.1, s["key_initiatives"], size=13, color=NAVY_DK,
                valign=MSO_ANCHOR.MIDDLE, shrink=True)
    footer(sl, d["company_name"])


def add_org_structure(prs, d):
    if not d.get("org_structure"):
        return
    sl = new_slide(prs)
    slide_title(sl, "Organization Structure")
    add_card(sl, 0.5, 1.25, PW - 1.0, PH - 1.85)
    add_textbox(sl, 0.8, 1.4, PW - 1.6, PH - 2.1, d["org_structure"], size=13, color=NAVY_DK,
                valign=MSO_ANCHOR.MIDDLE, shrink=True)
    footer(sl, d["company_name"])


def initials(name):
    parts = (name or "?").split()
    return "".join(w[0] for w in parts[:2]).upper() or "?"


def add_people_slide(prs, d, title, people):
    if not people:
        return
    sl = new_slide(prs)
    slide_title(sl, title)
    rows = [people[i:i + 3] for i in range(0, len(people), 3)]
    row_h = (PH - 1.5) / len(rows)
    for ri, row in enumerate(rows):
        col_w = (PW - 1.0 - 0.3 * (len(row) - 1)) / len(row)
        for ci, p in enumerate(row):
            x = 0.5 + ci * (col_w + 0.3)
            y = 1.25 + ri * row_h
            h = row_h - 0.25
            add_card(sl, x, y, col_w, h)
            photo = p.get("photo_path")
            if photo and os.path.exists(photo):
                add_circle_image(sl, photo, x + 0.18, y + 0.18, 0.65, 0.65)
            else:
                add_ellipse(sl, x + 0.18, y + 0.18, 0.65, 0.65, ACCENT)
                add_textbox(sl, x + 0.18, y + 0.18, 0.65, 0.65, initials(p.get("name")), size=16, bold=True,
                            color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            add_textbox(sl, x + 0.95, y + 0.15, col_w - 1.1, 0.35, p.get("name") or "", size=11.5, bold=True, color=NAVY, shrink=True)
            add_textbox(sl, x + 0.95, y + 0.48, col_w - 1.1, 0.45, p.get("designation") or "", size=9, color=SLATE, shrink=True)
            if p.get("brief"):
                add_textbox(sl, x + 0.18, y + 0.95, col_w - 0.36, h - 1.15, p["brief"], size=8.5, color=NAVY_DK, shrink=True)
            if p.get("linkedin_url"):
                add_textbox(sl, x + 0.18, y + h - 0.32, col_w - 0.36, 0.28, "LinkedIn ↗", size=8.5, color=ACCENT,
                            hyperlink=p["linkedin_url"])
    footer(sl, d["company_name"])


def add_swot_overview(prs, d):
    sw = d["swot"]
    quads = [("Strengths", sw.get("strengths") or [], GOOD), ("Weaknesses", sw.get("weaknesses") or [], BAD),
             ("Opportunities", sw.get("opportunities") or [], ACCENT), ("Threats", sw.get("threats") or [], WARN)]
    if not any(arr for _, arr, _ in quads):
        return
    sl = new_slide(prs)
    slide_title(sl, "SWOT Analysis")
    col_w = (PW - 1.0 - 0.3 * 3) / 4
    for i, (label, arr, color) in enumerate(quads):
        x = 0.5 + i * (col_w + 0.3)
        add_card(sl, x, 1.25, col_w, PH - 1.85, fill=LIGHT_BG)
        add_rect(sl, x, 1.25, col_w, 0.5, color)
        add_textbox(sl, x, 1.25, col_w, 0.5, label.upper(), size=12, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        pts = [r.get("point") for r in arr[:5] if r.get("point")]
        if pts:
            add_bullets(sl, x + 0.15, 1.9, col_w - 0.3, PH - 2.5, pts, size=9.5, color=NAVY_DK)
    footer(sl, d["company_name"])


def add_swot_detail(prs, d, title, arr, color):
    if not arr:
        return
    sl = new_slide(prs)
    slide_title(sl, f"SWOT – {title}")
    y = 1.25
    row_h = min(1.35, (PH - 1.85) / len(arr))
    for item in arr[:5]:
        add_card(sl, 0.5, y, PW - 1.0, row_h - 0.15)
        add_ellipse(sl, 0.7, y + 0.16, 0.14, 0.14, color)
        add_textbox(sl, 0.98, y + 0.08, PW - 1.73, 0.35, item.get("point") or "", size=12, bold=True, color=NAVY, shrink=True)
        body = item.get("detail") or item.get("evidence")
        if body:
            add_textbox(sl, 0.98, y + 0.42, PW - 1.73, row_h - 0.55, body, size=9.5, color=SLATE, shrink=True)
        y += row_h
    footer(sl, d["company_name"])


def add_financials_annual(prs, d):
    f = d["financials_annual"]
    has_current = any((f[k] or {}).get("current") is not None for k in ("revenue", "operating_income", "net_income"))
    if not has_current:
        return
    sl = new_slide(prs)
    slide_title(sl, "Financials – Annual")

    y_labels = ["2 yrs prior", "Prior year", f.get("financial_year") or "Current year"]
    raw_series = [
        ("Revenue ($)", [f["revenue"].get("two_prior"), f["revenue"].get("previous"), f["revenue"].get("current")]),
        ("Operating Income ($)", [f["operating_income"].get("two_prior"), f["operating_income"].get("previous"), f["operating_income"].get("current")]),
        ("Net Income ($)", [f["net_income"].get("two_prior"), f["net_income"].get("previous"), f["net_income"].get("current")]),
    ]
    series = [(name, [(v / 1e9 if v is not None else 0) for v in vals])
              for name, vals in raw_series if any(v is not None for v in vals)]

    if series:
        chart_data = CategoryChartData()
        chart_data.categories = y_labels
        for name, vals in series:
            chart_data.add_series(name, vals)
        gframe = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, IN(0.5), IN(1.25), IN(8.0), IN(5.6), chart_data)
        chart = gframe.chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "Revenue / Operating Income / Net Income ($B)"
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        palette = [NAVY, ACCENT, RGBColor(0x9D, 0xB4, 0xC9)]
        for i, s in enumerate(chart.plots[0].series):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = palette[i % len(palette)]
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.color.rgb = SLATE
        cat_axis.has_major_gridlines = False
        val_axis = chart.value_axis
        val_axis.tick_labels.font.color.rgb = SLATE
        val_axis.major_gridlines.format.line.color.rgb = LINE_GREY

    hi = f["highlights"]
    rows = [(l, v) for l, v in [("EBITDA", hi.get("ebitda")), ("Gross Profit", hi.get("gross_profit")),
                                  ("EPS", hi.get("eps")), ("Revenue Growth", hi.get("revenue_growth")),
                                  ("Operating Margin", hi.get("operating_margin")), ("Net Margin", hi.get("net_margin")),
                                  ("Free Cash Flow", hi.get("free_cash_flow")), ("ROE", hi.get("roe"))]
            if v is not None and v != ""]
    if rows:
        add_textbox(sl, 8.85, 1.25, 3.9, 0.35, "Highlights", size=13, bold=True, color=NAVY)
        y = 1.65
        for label, val in rows:
            add_textbox(sl, 8.85, y, 2.3, 0.35, label, size=10, color=SLATE)
            add_textbox(sl, 11.1, y, 1.6, 0.35, str(val), size=10, bold=True, color=NAVY_DK, align=PP_ALIGN.RIGHT)
            y += 0.4
    footer(sl, d["company_name"])


def add_financials_quarterly(prs, d):
    q = d["financials_quarterly"]
    rows = [(l, v) for l, v in [("Quarter", q.get("quarter")), ("Revenue", q.get("revenue")),
                                  ("Operating Income", q.get("operating_income")), ("Net Income", q.get("net_income")),
                                  ("EBITDA", q.get("ebitda")), ("EPS", q.get("eps")), ("Cash Flow", q.get("cash_flow")),
                                  ("Revenue Growth (YoY)", q.get("revenue_growth")),
                                  ("Net Income Growth (YoY)", q.get("net_income_growth"))]
            if v is not None and v != ""]
    if not rows:
        return
    sl = new_slide(prs)
    slide_title(sl, "Financials – Current Quarter")
    per_col = -(-len(rows) // 2)  # ceil
    cols = [rows[:per_col], rows[per_col:]]
    for ci, col in enumerate(cols):
        x = 0.5 + ci * ((PW - 1.0) / 2 + 0.1)
        y = 1.35
        for label, val in col:
            add_card(sl, x, y, (PW - 1.2) / 2, 0.85)
            add_textbox(sl, x + 0.15, y + 0.1, (PW - 1.2) / 2 - 0.3, 0.3, label.upper(), size=9, bold=True, color=SLATE)
            add_textbox(sl, x + 0.15, y + 0.38, (PW - 1.2) / 2 - 0.3, 0.4, str(val), size=14, bold=True, color=NAVY)
            y += 1.0
    footer(sl, d["company_name"])


def add_card_list_slide(prs, d, title, items, fields):
    """fields: list of dicts {key, label, primary?}"""
    if not items:
        return
    sl = new_slide(prs)
    slide_title(sl, title)
    cols = 1 if len(items) <= 2 else 2
    rows = -(-len(items) // cols)
    col_w = (PW - 1.0 - 0.25 * (cols - 1)) / cols
    row_h = min(1.6, (PH - 1.85 - 0.2 * (rows - 1)) / rows)
    shown = items[:cols * min(rows, 6)]
    primary = next((f for f in fields if f.get("primary")), None)
    for i, item in enumerate(shown):
        c, r = i % cols, i // cols
        x = 0.5 + c * (col_w + 0.25)
        y = 1.25 + r * (row_h + 0.2)
        add_card(sl, x, y, col_w, row_h)
        tx, ty = x + 0.15, y + 0.1
        if primary and item.get(primary["key"]):
            add_textbox(sl, tx, ty, col_w - 0.3, 0.32, item[primary["key"]], size=12, bold=True, color=NAVY, shrink=True)
            ty += 0.36
        rest = [f for f in fields if not f.get("primary") and item.get(f["key"])]
        body = "  |  ".join(f"{f['label']}: {item[f['key']]}" for f in rest)
        if body:
            add_textbox(sl, tx, ty, col_w - 0.3, row_h - (ty - y) - 0.1, body, size=9, color=SLATE, shrink=True)
    footer(sl, d["company_name"])


def add_table_slide(prs, d, title, headers, rows, col_ws):
    if not rows:
        return
    sl = new_slide(prs)
    slide_title(sl, title)
    n_rows, n_cols = len(rows) + 1, len(headers)
    tbl_shape = sl.shapes.add_table(n_rows, n_cols, IN(0.5), IN(1.25), IN(PW - 1.0), IN(0.4 * n_rows))
    table = tbl_shape.table
    for ci, w in enumerate(col_ws):
        table.columns[ci].width = IN(w)
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(10.5)
        p.runs[0].font.color.rgb = WHITE
        p.runs[0].font.name = FONT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val) if val not in (None, "") else "—"
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(9.5)
            p.runs[0].font.color.rgb = NAVY_DK
            p.runs[0].font.name = FONT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    footer(sl, d["company_name"])


def add_news(prs, d):
    groups = d["news"]
    cats = [(k, v) for k, v in groups.items() if v]
    if not cats:
        return
    sl = new_slide(prs)
    slide_title(sl, "Latest News")
    col_w = (PW - 1.0 - 0.25 * (len(cats) - 1)) / len(cats)
    for ci, (cat, items) in enumerate(cats):
        x = 0.5 + ci * (col_w + 0.25)
        add_card(sl, x, 1.2, col_w, 0.45, fill=NAVY)
        add_textbox(sl, x, 1.2, col_w, 0.45, cat, size=11, bold=True, color=WHITE,
                    align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        y = 1.8
        item_h = (PH - 2.35) / max(len(items), 1)
        for n in items:
            h = min(item_h, 1.5)
            add_card(sl, x, y, col_w, h - 0.12)
            if n.get("date"):
                add_textbox(sl, x + 0.12, y + 0.06, col_w - 0.24, 0.22, str(n["date"])[:10], size=8, color=SLATE)
            add_textbox(sl, x + 0.12, y + 0.26, col_w - 0.24, 0.4, n.get("title") or "", size=9.5, bold=True, color=NAVY, shrink=True)
            if n.get("summary"):
                add_textbox(sl, x + 0.12, y + 0.62, col_w - 0.24, h - 0.85, n["summary"], size=8, color=SLATE, shrink=True)
            if n.get("url"):
                add_textbox(sl, x + 0.12, y + h - 0.28, col_w - 0.24, 0.2, "Read more ↗", size=7.5, color=ACCENT, hyperlink=n["url"])
            y += h
    footer(sl, d["company_name"])


def add_it_spending(prs, d):
    if not d.get("it_spending"):
        return
    sl = new_slide(prs)
    slide_title(sl, "IT Spending")
    add_card(sl, 0.5, 1.25, PW - 1.0, PH - 1.85)
    add_textbox(sl, 0.8, 1.4, PW - 1.6, PH - 2.1, d["it_spending"], size=13, color=NAVY_DK,
                valign=MSO_ANCHOR.MIDDLE, shrink=True)
    footer(sl, d["company_name"])


def add_tech_initiatives(prs, d):
    items = d["tech_initiatives"]
    if not items:
        return
    sl = new_slide(prs)
    slide_title(sl, "Technology Initiatives")
    cols = 2
    rows = -(-len(items) // cols)
    col_w = (PW - 1.0 - 0.25) / cols
    row_h = min(1.5, (PH - 1.85 - 0.2 * (rows - 1)) / rows)
    shown = items[:cols * min(rows, 6)]
    for i, item in enumerate(shown):
        c, r = i % cols, i // cols
        x = 0.5 + c * (col_w + 0.25)
        y = 1.25 + r * (row_h + 0.2)
        add_card(sl, x, y, col_w, row_h)
        add_textbox(sl, x + 0.15, y + 0.08, col_w - 0.3, 0.25, item.get("date") or "", size=8.5, bold=True, color=ACCENT)
        add_textbox(sl, x + 0.15, y + 0.32, col_w - 0.3, 0.35, item.get("title") or "", size=11.5, bold=True, color=NAVY, shrink=True)
        if item.get("details"):
            add_textbox(sl, x + 0.15, y + 0.68, col_w - 0.3, row_h - 0.8, item["details"], size=9, color=SLATE, shrink=True)
    footer(sl, d["company_name"])


def add_technologies_in_use(prs, d):
    items = d["technologies_in_use"]
    if not items:
        return
    per_slide = 12
    for p_start in range(0, len(items), per_slide):
        chunk = items[p_start:p_start + per_slide]
        sl = new_slide(prs)
        title = f"Technologies in Use ({p_start // per_slide + 1})" if len(items) > per_slide else "Technologies in Use"
        slide_title(sl, title)
        cols = 3
        rows = -(-len(chunk) // cols)
        col_w = (PW - 1.0 - 0.2 * (cols - 1)) / cols
        row_h = min(1.1, (PH - 1.85 - 0.15 * (rows - 1)) / rows)
        for i, t in enumerate(chunk):
            c, r = i % cols, i // cols
            x = 0.5 + c * (col_w + 0.2)
            y = 1.25 + r * (row_h + 0.15)
            add_card(sl, x, y, col_w, row_h)
            add_textbox(sl, x + 0.12, y + 0.08, col_w - 0.24, 0.3, t.get("technology") or "", size=11, bold=True, color=NAVY, shrink=True)
            if t.get("category"):
                add_textbox(sl, x + 0.12, y + 0.36, col_w - 0.24, 0.22, t["category"], size=8, color=ACCENT)
            if t.get("brief"):
                add_textbox(sl, x + 0.12, y + 0.58, col_w - 0.24, row_h - 0.68, t["brief"], size=8, color=SLATE, shrink=True)
        footer(sl, d["company_name"])


def add_industry_indicators(prs, d):
    ind = d["industry_indicators"]
    rows = [(l, v) for l, v in [("Demand Trend", ind.get("demand_trend")), ("Technology Adoption", ind.get("technology_adoption")),
                                  ("Spending Trend", ind.get("spending_trend")), ("Regulatory Environment", ind.get("regulatory_environment")),
                                  ("Competitive Intensity", ind.get("competitive_intensity")), ("Outlook", ind.get("outlook"))]
            if v]
    if not ind.get("growth_rating") and not rows:
        return
    sl = new_slide(prs)
    slide_title(sl, "Industry Indicators")

    if ind.get("growth_rating"):
        levels = ["Low", "Medium", "High"]
        gr_lower = ind["growth_rating"].lower()
        active = next((i for i, l in enumerate(levels) if l.lower() in gr_lower), -1)
        add_textbox(sl, 0.5, 1.3, 3, 0.35, "Growth Rating", size=12, bold=True, color=NAVY)
        for i, lvl in enumerate(levels):
            x = 0.5 + i * 2.0
            is_active = (i == active) or (active == -1 and lvl.lower() == "medium")
            add_card(sl, x, 1.75, 1.8, 0.6, fill=(ACCENT if is_active else LIGHT_BG))
            add_textbox(sl, x, 1.75, 1.8, 0.6, lvl, size=11, bold=True, color=(WHITE if is_active else SLATE),
                        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    y = 2.75
    for label, val in rows:
        add_card(sl, 0.5, y, PW - 1.0, 0.55)
        add_textbox(sl, 0.7, y, 3.5, 0.55, label, size=11, bold=True, color=SLATE, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(sl, 4.3, y, PW - 5.0, 0.55, val, size=11, color=NAVY_DK, valign=MSO_ANCHOR.MIDDLE, shrink=True)
        y += 0.65
    footer(sl, d["company_name"])


def add_industry_forecast(prs, d):
    f = d["industry_forecast"]
    items = [(l, v) for l, v in [("Industry", f.get("industry")), ("Competitive Landscape", f.get("competitive_landscape")),
                                   ("Forecast", f.get("forecast")), ("Growth Drivers", f.get("growth_drivers"))] if v]
    if not items:
        return
    sl = new_slide(prs)
    slide_title(sl, "Industry Forecast & Competitive Landscape")
    y = 1.25
    total_len = sum(len(v) for _, v in items) or 1
    for label, val in items:
        h = max(0.9, (PH - 1.85) * (len(val) / total_len))
        add_card(sl, 0.5, y, PW - 1.0, h - 0.15)
        add_textbox(sl, 0.75, y + 0.1, PW - 1.5, 0.3, label, size=12, bold=True, color=ACCENT)
        add_textbox(sl, 0.75, y + 0.42, PW - 1.5, h - 0.6, val, size=9.5, color=NAVY_DK, shrink=True)
        y += h
    footer(sl, d["company_name"])


def add_thank_you(prs, d):
    sl = new_slide(prs, NAVY_DK)
    add_textbox(sl, 0.8, 3.0, PW - 1.6, 1.0, "Thank You", size=40, bold=True, color=WHITE)
    add_textbox(sl, 0.8, 3.9, PW - 1.6, 0.5, f'{d["company_name"]} — Company Intelligence', size=14, color=RGBColor(0xAF, 0xC6, 0xE3))


# ---------------------------------------------------------------------------
def build(d, out_path):
    prs = Presentation()
    prs.slide_width = IN(PW)
    prs.slide_height = IN(PH)

    add_title_slide(prs, d)                                                          # 1
    add_overview(prs, d)                                                             # 2
    add_mission_vision(prs, d)                                                       # 3
    add_geo(prs, d)                                                                  # 4
    add_segments(prs, d)                                                             # 5
    add_sustainability_strategy(prs, d)                                              # 6
    add_sustainability_initiatives(prs, d)                                           # 7
    add_org_structure(prs, d)                                                        # 8
    add_people_slide(prs, d, "Leadership Team", d["leadership"])                     # 9
    add_people_slide(prs, d, "Organization Structure – Technology Team", d.get("technology_team") or [])
    add_separator(prs, "SWOT Analysis")
    add_swot_overview(prs, d)                                                        # 10
    add_swot_detail(prs, d, "Strengths", d["swot"].get("strengths") or [], GOOD)      # 11
    add_swot_detail(prs, d, "Weaknesses", d["swot"].get("weaknesses") or [], BAD)     # 12
    add_swot_detail(prs, d, "Opportunities", d["swot"].get("opportunities") or [], ACCENT)  # 13
    add_swot_detail(prs, d, "Threats", d["swot"].get("threats") or [], WARN)          # 14
    add_separator(prs, "Financials")
    add_financials_annual(prs, d)                                                    # 15
    add_financials_quarterly(prs, d)                                                 # 16
    add_card_list_slide(prs, d, "Key Acquisitions", d["acquisitions"],
                         [{"key": "company_name", "label": "Company", "primary": True},
                          {"key": "year", "label": "Year"}, {"key": "value", "label": "Value"},
                          {"key": "brief", "label": "Brief"}])                        # 17
    add_table_slide(prs, d, "Key Competitors",
                     ["Company", "Revenue", "Employees", "Market Cap", "ICT Budget"],
                     [[c.get("name"), c.get("revenue"), c.get("employees"), c.get("market_cap"), c.get("ict_budget")]
                      for c in d["competitors"]],
                     [3.2, 2.4, 2.4, 2.4, 2.13])                                      # 18
    add_card_list_slide(prs, d, "Awards and Accolades", d["awards"],
                         [{"key": "award", "label": "Award", "primary": True},
                          {"key": "date", "label": "Date"}, {"key": "brief", "label": "Brief"}])  # 19
    add_card_list_slide(prs, d, "Business Challenges", d["challenges"],
                         [{"key": "challenge", "label": "Challenge", "primary": True},
                          {"key": "impact", "label": "Impact"}, {"key": "brief", "label": "Brief"}])  # 20
    add_separator(prs, "Market & Technology")
    add_news(prs, d)                                                                 # 21
    add_it_spending(prs, d)                                                          # 22
    add_table_slide(prs, d, "Deals", ["Vendor", "Start Date", "End Date", "Contract Details"],
                     [[dl.get("vendor"), dl.get("start_date"), dl.get("end_date"), dl.get("contract_details")]
                      for dl in d["deals"]],
                     [2.8, 1.6, 1.6, 6.33])                                           # 23
    add_tech_initiatives(prs, d)                                                     # 24
    add_technologies_in_use(prs, d)                                                  # 25
    add_separator(prs, "Industry Outlook")
    add_industry_indicators(prs, d)                                                  # 26
    add_industry_forecast(prs, d)                                                    # 27
    add_thank_you(prs, d)                                                            # 28

    prs.save(out_path)
    print(f"wrote {out_path}")


def main():
    if len(sys.argv) != 3:
        print("usage: build_ppt.py <slide_data.json> <out.pptx>", file=sys.stderr)
        sys.exit(1)
    with open(sys.argv[1], encoding="utf-8") as f:
        d = json.load(f)
    build(d, sys.argv[2])


if __name__ == "__main__":
    main()
